"""Rute pembangkit dokumen — Slide Presentasi (PPTX) & Proposal (DOCX).

Kedua dokumen adalah materi PRD/penawaran AMAN (Aplikasi Manajemen Aset
Negara) yang bisa diunduh dari halaman Info (login → #fitur). ISINYA WAJIB
MENGIKUTI VERSI APLIKASI TERKINI — sumber kebenarannya halaman PRD
(`frontend/src/pages/InfoPage.jsx`) + README; angka RAB & skema lisensi di
sini harus selalu SAMA dengan seksi "Harga, Lisensi & RAB" halaman itu.
"""
import io
import logging
from datetime import datetime

from fastapi import APIRouter, Request
from fastapi.responses import StreamingResponse

from shared_utils import limiter

logger = logging.getLogger(__name__)
documents_router = APIRouter()

# ── Identitas & angka kunci (satu tempat agar PPT dan DOCX tak saling
#    berbeda; selaraskan dengan InfoPage saat aplikasi berkembang) ──
NAMA_APLIKASI = "AMAN"
NAMA_PANJANG = "Aplikasi Manajemen Aset Negara"
VERSI = "v2.6"
STAT_UTAMA = [("12", "Tahap Siklus BMN"), ("31", "Halaman Aplikasi"),
              ("15+", "Laporan Resmi PDF"), ("45", "Kode Transaksi SAKTI")]
TIER_LISENSI = [
    ("Lisensi Satker Tunggal", "Rp 55 juta / tahun",
     "Satu satuan kerja, pengguna tidak dibatasi", [
         "Semua fitur (offline, real-time, 15+ laporan, pengesahan)",
         "Update versi & perbaikan bug",
         "Dukungan teknis jam kerja",
         "Instalasi di server satker (on-premise) atau VPS sendiri"]),
    ("Multi-Satker / Instansi", "Rp 175 juta / tahun",
     "Hingga 10 satker dalam satu instansi", [
         "Semua benefit Satker Tunggal",
         "Kop surat & pengaturan laporan per satker",
         "Dukungan prioritas + pendampingan implementasi",
         "Pelatihan pengguna (2 sesi / tahun)"]),
    ("Perpetual + Source Code", "Rp 500 juta sekali bayar",
     "Hak pakai selamanya + kendali penuh", [
         "Lisensi permanen tanpa biaya tahunan wajib",
         "Serah terima source code lengkap + dokumentasi",
         "Transfer knowledge ke tim TI instansi",
         "Dukungan teknis 12 bulan pertama"]),
]
# RAB referensi "membangun sendiri" — identik dengan RABTable di InfoPage.
RAB_BARIS = [
    ("A", "BIAYA PENGEMBANGAN", "", "", "", ""),
    ("1", "Full-stack Development", "1", "Paket", "85.000.000", "85.000.000"),
    ("2", "UI/UX Design & Prototyping", "1", "Paket", "15.000.000", "15.000.000"),
    ("3", "Quality Assurance & Testing", "1", "Paket", "10.000.000", "10.000.000"),
    ("4", "Dokumentasi & User Manual", "1", "Paket", "5.000.000", "5.000.000"),
    ("5", "Training Pengguna (2 sesi)", "2", "Sesi", "5.000.000", "10.000.000"),
    ("", "Subtotal Pengembangan", "", "", "", "125.000.000"),
    ("B", "BIAYA INFRASTRUKTUR (1 TAHUN)", "", "", "", ""),
    ("6", "Cloud Server (VPS)", "12", "Bulan", "1.500.000", "18.000.000"),
    ("7", "Domain & SSL Certificate", "1", "Tahun", "500.000", "500.000"),
    ("8", "MongoDB Atlas (M10 Cluster)", "12", "Bulan", "800.000", "9.600.000"),
    ("9", "Backup Storage (100 GB)", "12", "Bulan", "200.000", "2.400.000"),
    ("", "Subtotal Infrastruktur", "", "", "", "30.500.000"),
    ("C", "BIAYA OPERASIONAL (1 TAHUN)", "", "", "", ""),
    ("10", "Maintenance & Bug Fix", "12", "Bulan", "2.000.000", "24.000.000"),
    ("11", "Technical Support", "12", "Bulan", "1.000.000", "12.000.000"),
    ("", "Subtotal Operasional", "", "", "", "36.000.000"),
]
RAB_TOTAL_TH1 = "191.500.000"
RAB_TAHUNAN = "66.500.000"
RAB_TOTAL_3TH = "324.500.000"


# ═══════════════════════════════════════════════════════════════════
# GENERATOR SLIDE PRESENTASI (PPTX)
# ═══════════════════════════════════════════════════════════════════
@documents_router.get("/documents/ppt")
@limiter.limit("5/minute")
async def generate_ppt(request: Request):
    """Slide presentasi PRD AMAN — 10 slide sesuai versi aplikasi terkini."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Palet warna (selaras nuansa teal halaman PRD)
    TEAL = RGBColor(0x0D, 0x94, 0x88)
    DARK = RGBColor(0x0F, 0x17, 0x2A)
    CARD = RGBColor(0x1E, 0x29, 0x3B)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x94, 0xA3, 0xB8)
    ACCENT = RGBColor(0x5E, 0xEA, 0xD4)
    BLUE = RGBColor(0x38, 0xBD, 0xF8)
    GREEN = RGBColor(0x22, 0xC5, 0x5E)
    ORANGE = RGBColor(0xF9, 0x73, 0x16)
    RED = RGBColor(0xEF, 0x44, 0x44)
    PURPLE = RGBColor(0xA7, 0x55, 0xF5)
    PINK = RGBColor(0xEC, 0x48, 0x99)

    def add_bg(slide, color=DARK):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape_bg(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width,
                                       height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = alignment
        return box

    def kepala_slide(slide, nomor, seksi, judul):
        add_bg(slide, DARK)
        add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333),
                     Inches(0.08), TEAL)
        add_text_box(slide, Inches(0.8), Inches(0.45), Inches(8), Inches(0.5),
                     f"{nomor}  {seksi}", 14, ACCENT, True)
        add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
                     judul, 30, WHITE, True)

    # ── SLIDE 1: Sampul ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), TEAL)
    add_text_box(slide, Inches(1), Inches(1.3), Inches(9), Inches(0.6),
                 f"{NAMA_APLIKASI} — {NAMA_PANJANG}", 16, ACCENT, True)
    add_text_box(slide, Inches(1), Inches(1.95), Inches(11), Inches(1.6),
                 "Platform Siklus Penuh Pengelolaan\nBarang Milik Negara (BMN)",
                 40, WHITE, True)
    add_text_box(slide, Inches(1), Inches(3.95), Inches(11), Inches(1.1),
                 "PWA offline-first · sinkron SIMAN V2 · persediaan standar "
                 "SAKTI · tanda tangan elektronik · peta GIS\nSelaras "
                 "SE-17/MK.1/2024, PP 27/2014 jo PP 28/2020, PMK 181/PMK.06/"
                 "2016, PMK 234/PMK.05/2020", 15, GRAY)
    add_text_box(slide, Inches(1), Inches(6.35), Inches(6), Inches(0.4),
                 f"Product Requirements Document {VERSI} | "
                 f"{datetime.now().strftime('%B %Y')}", 12, GRAY)
    for i, (val, label) in enumerate(STAT_UTAMA):
        x = Inches(7.5) + Inches(i * 1.4)
        add_shape_bg(slide, x, Inches(5.35), Inches(1.25), Inches(1.25), CARD)
        add_text_box(slide, x, Inches(5.45), Inches(1.25), Inches(0.5), val,
                     19, ACCENT, True, PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(6.0), Inches(1.25), Inches(0.55), label,
                     8, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 2: Latar Belakang ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kepala_slide(slide, "01", "LATAR BELAKANG", "Mengapa Sistem Ini Dibutuhkan?")
    masalah = [
        ("Pencatatan Manual", "Spreadsheet rawan salah input, duplikasi, dan "
         "sulit diaudit; rekap laporan memakan waktu berhari-hari."),
        ("Sinyal Lapangan Terbatas", "Inventarisasi fisik sering di lokasi "
         "tanpa koneksi — sistem daring biasa berhenti bekerja."),
        ("Kepatuhan Regulasi", "SE-17/MK.1/2024 menuntut klasifikasi rinci "
         "aset tidak ditemukan; SAKTI menuntut kode transaksi & FIFO "
         "persediaan yang persis."),
        ("Data Tersebar", "Foto, dokumen, BAST, nilai, dan riwayat aset "
         "terpisah-pisah sehingga sulit diverifikasi dan direkonsiliasi."),
    ]
    for i, (judul, isi) in enumerate(masalah):
        y = Inches(2.0) + Inches(i * 1.25)
        add_shape_bg(slide, Inches(0.8), y, Inches(0.06), Inches(1.0),
                     [BLUE, ORANGE, RED, GREEN][i])
        add_text_box(slide, Inches(1.1), y, Inches(5.4), Inches(0.4), judul,
                     16, WHITE, True)
        add_text_box(slide, Inches(1.1), y + Inches(0.42), Inches(5.4),
                     Inches(0.8), isi, 11, GRAY)
    add_shape_bg(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(5.1),
                 CARD)
    add_text_box(slide, Inches(7.3), Inches(2.15), Inches(5), Inches(0.4),
                 "DASAR HUKUM", 12, ACCENT, True)
    regulasi = [
        "SE-17/MK.1/2024 — Pelaksanaan\nInventarisasi BMN",
        "PP 27/2014 jo PP 28/2020 —\nPengelolaan BMN/D",
        "PMK 181/PMK.06/2016 —\nPenatausahaan BMN",
        "PMK 234/PMK.05/2020 — Kebijakan\nAkuntansi (FIFO Persediaan)",
    ]
    for i, reg in enumerate(regulasi):
        y = Inches(2.75) + Inches(i * 1.02)
        add_shape_bg(slide, Inches(7.4), y, Inches(0.35), Inches(0.35), TEAL)
        add_text_box(slide, Inches(7.44), y - Inches(0.02), Inches(0.35),
                     Inches(0.35), str(i + 1), 12, WHITE, True,
                     PP_ALIGN.CENTER)
        add_text_box(slide, Inches(8.0), y - Inches(0.05), Inches(4.3),
                     Inches(0.9), reg, 12, WHITE)

    # ── SLIDE 3: Arsitektur ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kepala_slide(slide, "02", "OVERVIEW SISTEM", "Arsitektur & Teknologi")
    teknologi = [
        ("Frontend", "React 19 + Tailwind CSS\nshadcn/ui + Leaflet (peta)\n"
         "PWA offline-first: IndexedDB\n+ Service Worker + antrean sinkron",
         BLUE),
        ("Backend", "Python FastAPI (async) + Motor\nWebSocket lintas worker\n"
         "OCC (version/If-Match)\n+ Idempotency-Key, JWT + OTP", GREEN),
        ("Database", "MongoDB 7 + GridFS\n(foto, dokumen, BAST)\nIndeks "
         "kunci + capped\ncollection ws_events", ORANGE),
        ("Infrastruktur", "VPS: Nginx + Supervisor\nCI/CD GitHub Actions — "
         "uji tiap\nPR, auto-deploy saat merge\nRedis & Meilisearch "
         "(opsional)", PURPLE),
    ]
    for i, (judul, isi, warna) in enumerate(teknologi):
        x = Inches(0.8) + Inches(i * 3.1)
        add_shape_bg(slide, x, Inches(2.0), Inches(2.85), Inches(3.0), CARD)
        add_shape_bg(slide, x, Inches(2.0), Inches(2.85), Inches(0.06), warna)
        add_text_box(slide, x + Inches(0.25), Inches(2.3), Inches(2.35),
                     Inches(0.4), judul, 16, warna, True)
        add_text_box(slide, x + Inches(0.25), Inches(2.8), Inches(2.45),
                     Inches(2.1), isi, 11, GRAY)
    add_text_box(slide, Inches(0.8), Inches(5.35), Inches(12), Inches(0.4),
                 "ALUR DATA", 12, ACCENT, True)
    alur = ["Peramban Pengguna", "React SPA (offline-ready)", "Nginx",
            "FastAPI multi-worker", "MongoDB + GridFS"]
    for i, item in enumerate(alur):
        x = Inches(0.6) + Inches(i * 2.55)
        add_shape_bg(slide, x, Inches(5.85), Inches(2.15), Inches(0.7), CARD)
        add_text_box(slide, x, Inches(5.92), Inches(2.15), Inches(0.6), item,
                     11, WHITE, False, PP_ALIGN.CENTER)
        if i < len(alur) - 1:
            add_text_box(slide, x + Inches(2.13), Inches(5.9), Inches(0.45),
                         Inches(0.6), "→", 18, ACCENT, True, PP_ALIGN.CENTER)

    # ── SLIDE 4: Siklus Penuh BMN ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kepala_slide(slide, "03", "CAKUPAN", "Siklus Penuh BMN — 12 Tahap + Modul Pendukung")
    tahap = ["Perencanaan", "Penganggaran", "Pengadaan", "Penggunaan",
             "Pemanfaatan", "Pengamanan & Pemeliharaan", "Penilaian",
             "Pemindahtanganan", "Pemusnahan", "Penghapusan",
             "Penatausahaan", "Wasdal"]
    for i, t in enumerate(tahap):
        kol, brs = i % 4, i // 4
        x = Inches(0.8) + Inches(kol * 3.05)
        y = Inches(1.95) + Inches(brs * 0.85)
        add_shape_bg(slide, x, y, Inches(2.85), Inches(0.7), CARD)
        add_shape_bg(slide, x, y, Inches(0.06), Inches(0.7), TEAL)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.08), Inches(2.6),
                     Inches(0.55), f"{i + 1:02d}  {t}", 12, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(4.75), Inches(12), Inches(0.4),
                 "PENATAUSAHAAN TIGA SERANGKAI + MODUL PENDUKUNG", 12, ACCENT,
                 True)
    pendukung = [
        ("Pembukuan", "DBKP intra/ekstra, Buku Barang\n(jurnal SIMAK/SAKTI), "
         "KIB A–F,\npenyusutan per semester", BLUE),
        ("Inventarisasi", "Kegiatan lapangan offline,\nkamera + QR + GPS, "
         "pengesahan\nberkekuatan dokumen", GREEN),
        ("Pelaporan", "LBKP/LBP semesteran-tahunan,\nrekonsiliasi XLSX, "
         "CaLBMN pra-isi,\narsip final terkunci", ORANGE),
        ("Persediaan SAKTI", "45 kode transaksi (M/K/P/H),\nFIFO per layer, "
         "opname + BAOF,\n7 laporan format SAKTI", PINK),
    ]
    for i, (judul, isi, warna) in enumerate(pendukung):
        x = Inches(0.8) + Inches(i * 3.1)
        add_shape_bg(slide, x, Inches(5.25), Inches(2.85), Inches(1.85), CARD)
        add_shape_bg(slide, x, Inches(5.25), Inches(2.85), Inches(0.05), warna)
        add_text_box(slide, x + Inches(0.2), Inches(5.42), Inches(2.5),
                     Inches(0.35), judul, 13, warna, True)
        add_text_box(slide, x + Inches(0.2), Inches(5.82), Inches(2.55),
                     Inches(1.2), isi, 10, GRAY)

    # ── SLIDE 5: Fitur Utama ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kepala_slide(slide, "04", "FITUR UTAMA", "Sorotan Kemampuan Terkini")
    fitur = [
        ("📶", "Offline Penuh", "Snapshot IndexedDB +\nantrean simpan persisten",
         ORANGE),
        ("📷", "Kamera Lapangan", "Watermark GPS/jam, flash,\nscan QR beruntun",
         GREEN),
        ("🗺️", "Peta GIS & Denah", "Leaflet + KML/KMZ/SHP,\ndenah dalam-gedung",
         BLUE),
        ("🔄", "Sinkron SIMAN V2", "Impor/ekspor Master Aset,\ntandai selisih per aset",
         ACCENT),
        ("📦", "Persediaan SAKTI", "45 kode + FIFO layer +\n7 laporan format resmi",
         PINK),
        ("📚", "Pembukuan & Nilai", "Jurnal Buku Barang, KIB,\nriwayat nilai + penyusutan",
         PURPLE),
        ("📑", "15+ Laporan Resmi", "DBHI/RHI/BAHI/SPTJM/LBP,\nkop resmi + batch ZIP",
         GREEN),
        ("✍️", "TTD Elektronik", "Link per penanda tangan,\nQR + hash verifikasi",
         RED),
        ("🏷️", "Stiker Label QR", "3 ukuran × A4/A3, hierarki\nterbaca + contoh berukuran",
         ORANGE),
        ("👥", "Master SDM", "Pegawai/pejabat/unit kerja,\nkartu UID e-KTP",
         BLUE),
        ("🛡️", "Multi-Satker Aman", "Isolasi data per satker,\naudit trail per field",
         PURPLE),
        ("💾", "Backup Otomatis", "Terjadwal harian + retensi,\npulihkan dari arsip",
         ACCENT),
    ]
    for i, (ikon, judul, isi, warna) in enumerate(fitur):
        kol, brs = i % 6, i // 6
        x = Inches(0.5) + Inches(kol * 2.1)
        y = Inches(1.85) + Inches(brs * 2.65)
        add_shape_bg(slide, x, y, Inches(1.95), Inches(2.35), CARD)
        add_shape_bg(slide, x, y, Inches(1.95), Inches(0.05), warna)
        add_text_box(slide, x, y + Inches(0.12), Inches(1.95), Inches(0.5),
                     ikon, 26, WHITE, False, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.68), Inches(1.75),
                     Inches(0.5), judul, 12, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(1.2), Inches(1.75),
                     Inches(1.05), isi, 9, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 6: Alur Kerja + Peran ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kepala_slide(slide, "05", "ALUR KERJA", "Workflow Inventarisasi & Peran Pengguna")
    langkah = [
        ("01", "Buat Kegiatan", "Nomor tiket INV-{tahun}-\n{seq}, satker & tim",
         BLUE),
        ("02", "Siapkan Data", "Impor SIMAN V2/CSV/XLSX\n46 kolom atau manual",
         GREEN),
        ("03", "Inventarisasi", "Lapangan offline: status,\nfoto, GPS, stiker, QR",
         ORANGE),
        ("04", "Verifikasi", "Klasifikasi SE-17, audit\ntrail, peta sebaran", RED),
        ("05", "Laporan", "15+ laporan resmi ber-kop\n+ batch unduh ZIP",
         PURPLE),
        ("06", "Sahkan & Kunci", "PDF ber-TTD → kegiatan\nterkunci + kartu riwayat",
         ACCENT),
    ]
    for i, (num, judul, isi, warna) in enumerate(langkah):
        x = Inches(0.5) + Inches(i * 2.1)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.72),
                                      Inches(1.95), Inches(0.6), Inches(0.6))
        circ.fill.solid()
        circ.fill.fore_color.rgb = warna
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(15)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if i < len(langkah) - 1:
            add_shape_bg(slide, x + Inches(1.38), Inches(2.2), Inches(1.4),
                         Inches(0.04), RGBColor(0x33, 0x44, 0x55))
        add_text_box(slide, x, Inches(2.7), Inches(2.0), Inches(0.4), judul,
                     13, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(3.1), Inches(2.0), Inches(1.0), isi, 10,
                     GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(0.4),
                 "PERAN PENGGUNA (RBAC + ISOLASI SATKER)", 12, ACCENT, True)
    peran = [
        ("Super Admin & Admin", "Kelola kegiatan/user/master, pengesahan & "
         "kunci, backup-restore-reset (khusus super admin), lintas satker "
         "hanya super admin", BLUE),
        ("Operator", "Input aset + foto + stiker + dokumen (termasuk "
         "offline), transaksi persediaan, scan QR, booking nomor surat",
         GREEN),
        ("Viewer", "Baca-saja (ditegakkan server): lihat data, unduh "
         "laporan, cetak kartu BMN & rekapitulasi", GRAY),
    ]
    for i, (nama, isi, warna) in enumerate(peran):
        x = Inches(0.8) + Inches(i * 4.2)
        add_shape_bg(slide, x, Inches(5.0), Inches(3.9), Inches(1.9), CARD)
        add_shape_bg(slide, x, Inches(5.0), Inches(0.06), Inches(1.9), warna)
        add_text_box(slide, x + Inches(0.25), Inches(5.2), Inches(3.4),
                     Inches(0.35), nama, 14, warna, True)
        add_text_box(slide, x + Inches(0.25), Inches(5.6), Inches(3.45),
                     Inches(1.2), isi, 10, GRAY)

    # ── SLIDE 7: Kepatuhan SE-17 & SAKTI ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kepala_slide(slide, "06", "KEPATUHAN", "Klasifikasi SE-17/MK.1/2024 & Standar SAKTI")
    status = [
        ("Ditemukan", "Aset terverifikasi keberadaannya di lokasi", GREEN, [
            "Kondisi: Baik / Rusak Ringan / Rusak Berat",
            "Foto aset wajib dilampirkan",
            "Stiker inventaris QR dipasang"]),
        ("Tidak Ditemukan", "Klasifikasi + sub-klasifikasi rinci", RED, [
            "Kesalahan Pencatatan (7 sub-klasifikasi)",
            "Tidak Ditemukan Lainnya (3 sub-klasifikasi)",
            "Uraian, kronologis & tindak lanjut"]),
        ("Status Lainnya", "Total 5 status inventarisasi didukung", GRAY, [
            "Belum Diinventarisasi (default impor, progres %)",
            "Berlebih — keterangan & asal-usul",
            "Sengketa — nomor perkara & pihak"]),
    ]
    for i, (nama, isi, warna, butir) in enumerate(status):
        x = Inches(0.5) + Inches(i * 4.2)
        add_shape_bg(slide, x, Inches(1.85), Inches(3.9), Inches(3.15), CARD)
        add_shape_bg(slide, x, Inches(1.85), Inches(3.9), Inches(0.06), warna)
        add_text_box(slide, x + Inches(0.3), Inches(2.1), Inches(3.3),
                     Inches(0.4), nama, 17, warna, True)
        add_text_box(slide, x + Inches(0.3), Inches(2.55), Inches(3.3),
                     Inches(0.5), isi, 11, GRAY)
        for j, b in enumerate(butir):
            add_text_box(slide, x + Inches(0.3), Inches(3.15) + Inches(j * 0.55),
                         Inches(3.4), Inches(0.5), f"• {b}", 10, WHITE)
    add_shape_bg(slide, Inches(0.5), Inches(5.25), Inches(12.3), Inches(1.7),
                 CARD)
    add_shape_bg(slide, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.05),
                 PINK)
    add_text_box(slide, Inches(0.8), Inches(5.42), Inches(11), Inches(0.35),
                 "PERSEDIAAN STANDAR SAKTI (PMK 234/PMK.05/2020)", 13, PINK,
                 True)
    add_text_box(slide, Inches(0.8), Inches(5.85), Inches(11.8), Inches(1.0),
                 "45 kode transaksi resmi (M01–M99 masuk, K01–K99 keluar, P01 "
                 "opname, H01–H03 hapus definitif ber-SK) · FIFO per layer · "
                 "Daftar Transaksi lintas barang · daftar Usang/Rusak/Tak "
                 "Dikuasai (bahan CaLK) · 7 laporan PDF persis format SAKTI",
                 12, GRAY)

    # ── SLIDE 8: Harga, Lisensi & RAB ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kepala_slide(slide, "07", "INVESTASI", "Skema Lisensi & Referensi RAB")
    for i, (nama, harga, tagline, _fitur) in enumerate(TIER_LISENSI):
        x = Inches(0.5) + Inches(i * 4.2)
        add_shape_bg(slide, x, Inches(1.9), Inches(3.9), Inches(2.5), CARD)
        add_shape_bg(slide, x, Inches(1.9), Inches(3.9), Inches(0.06),
                     [BLUE, TEAL, PURPLE][i])
        add_text_box(slide, x + Inches(0.3), Inches(2.1), Inches(3.3),
                     Inches(0.4), nama, 14, WHITE, True)
        add_text_box(slide, x + Inches(0.3), Inches(2.55), Inches(3.3),
                     Inches(0.5), harga, 18, [BLUE, ACCENT, PURPLE][i], True)
        add_text_box(slide, x + Inches(0.3), Inches(3.1), Inches(3.4),
                     Inches(0.5), tagline, 11, GRAY)
        add_text_box(slide, x + Inches(0.3), Inches(3.6), Inches(3.4),
                     Inches(0.75), "• " + "\n• ".join(_fitur[:2]), 9, GRAY)
    add_text_box(slide, Inches(0.8), Inches(4.7), Inches(12), Inches(0.35),
                 "REFERENSI RAB — MEMBANGUN & MENGOPERASIKAN SENDIRI", 12,
                 ACCENT, True)
    ringkas = [
        ("Tahun 1 (Dev + Infra + Ops)", f"Rp {RAB_TOTAL_TH1}"),
        ("Tahun 2 & 3 (Infra + Ops)", f"Rp {RAB_TAHUNAN} / tahun"),
        ("Total Investasi 3 Tahun", f"Rp {RAB_TOTAL_3TH}"),
    ]
    for i, (label, nilai) in enumerate(ringkas):
        x = Inches(0.8) + Inches(i * 4.2)
        add_shape_bg(slide, x, Inches(5.15), Inches(3.9), Inches(1.1), CARD)
        add_text_box(slide, x + Inches(0.25), Inches(5.3), Inches(3.4),
                     Inches(0.4), label, 11, GRAY)
        add_text_box(slide, x + Inches(0.25), Inches(5.7), Inches(3.4),
                     Inches(0.45), nilai, 16, ACCENT, True)
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.6),
                 "Harga acuan penawaran — nilai final mengikuti negosiasi, "
                 "lingkup kustomisasi, dan ketentuan pengadaan yang berlaku. "
                 "Tanpa biaya per-user.", 10, GRAY)

    # ── SLIDE 9: Timeline ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kepala_slide(slide, "08", "TIMELINE IMPLEMENTASI", "Jadwal Pengembangan & Pendampingan")
    timeline = [
        ("Bulan 1-2", "Analisis & Desain", "Requirement gathering\nUI/UX "
         "prototyping\nDesain database", BLUE),
        ("Bulan 3-5", "Pengembangan Core", "Backend API\nFrontend UI\n"
         "Database setup", GREEN),
        ("Bulan 6-7", "Integrasi & Testing", "Integrasi modul\nUAT & "
         "perbaikan\nPerformance tuning", ORANGE),
        ("Bulan 8", "Deployment", "Deploy produksi\nMigrasi data\nPelatihan "
         "pengguna", PURPLE),
        ("Bulan 9-12", "Maintenance", "Perbaikan & update\nPeningkatan fitur\n"
         "Dukungan teknis", ACCENT),
    ]
    add_shape_bg(slide, Inches(0.8), Inches(2.75), Inches(11.7), Inches(0.04),
                 RGBColor(0x33, 0x44, 0x55))
    for i, (periode, judul, isi, warna) in enumerate(timeline):
        x = Inches(0.5) + Inches(i * 2.5)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85),
                                      Inches(2.57), Inches(0.35), Inches(0.35))
        circ.fill.solid()
        circ.fill.fore_color.rgb = warna
        circ.line.fill.background()
        add_text_box(slide, x, Inches(2.0), Inches(2.3), Inches(0.4), periode,
                     11, warna, True, PP_ALIGN.CENTER)
        add_shape_bg(slide, x, Inches(3.3), Inches(2.3), Inches(2.6), CARD)
        add_shape_bg(slide, x, Inches(3.3), Inches(2.3), Inches(0.05), warna)
        add_text_box(slide, x + Inches(0.2), Inches(3.5), Inches(1.95),
                     Inches(0.45), judul, 13, WHITE, True)
        add_text_box(slide, x + Inches(0.2), Inches(4.05), Inches(1.95),
                     Inches(1.7), isi, 10, GRAY)

    # ── SLIDE 10: Penutup ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), TEAL)
    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
                 "Terima Kasih", 48, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
                 f"{NAMA_APLIKASI} — pengelolaan BMN satu pintu dari "
                 "perencanaan sampai penghapusan,\nselaras SIMAN V2 & SAKTI, "
                 "tetap bekerja penuh saat sinyal hilang", 18, GRAY, False,
                 PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
                 "Siap untuk demo & diskusi lebih lanjut", 16, ACCENT, False,
                 PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
                 f"© {datetime.now().year} {NAMA_APLIKASI} — {NAMA_PANJANG} | "
                 f"Product Requirements Document {VERSI}", 11, GRAY, False,
                 PP_ALIGN.CENTER)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument."
                   "presentationml.presentation",
        headers={"Content-Disposition":
                 "attachment; filename=AMAN_PRD_Presentasi.pptx"})


# ═══════════════════════════════════════════════════════════════════
# GENERATOR PROPOSAL (DOCX)
# ═══════════════════════════════════════════════════════════════════
@documents_router.get("/documents/proposal")
@limiter.limit("5/minute")
async def generate_proposal(request: Request):
    """Proposal DOCX BAB I–VI sesuai versi aplikasi terkini + RAB."""
    from docx import Document
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Cm
    from docx.shared import Pt as DPt
    from docx.shared import RGBColor as DRGBColor

    doc = Document()

    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = DPt(11)
    style.paragraph_format.space_after = DPt(6)
    style.paragraph_format.line_spacing = 1.15
    for section in doc.sections:
        section.top_margin = Cm(2.54)
        section.bottom_margin = Cm(2.54)
        section.left_margin = Cm(3)
        section.right_margin = Cm(2.54)

    def add_heading_styled(text, level=1):
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.color.rgb = DRGBColor(0x0D, 0x94, 0x88)
        return h

    def add_para(text, bold=False, italic=False, alignment=None,
                 space_after=6):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.bold = bold
        run.italic = italic
        run.font.size = DPt(11)
        run.font.name = 'Calibri'
        if alignment:
            p.alignment = alignment
        p.paragraph_format.space_after = DPt(space_after)
        return p

    def add_table_row(table, cells_data, bold=False, bg_color=None):
        row = table.add_row()
        for i, text in enumerate(cells_data):
            cell = row.cells[i]
            cell.text = ""
            run = cell.paragraphs[0].add_run(str(text))
            run.bold = bold
            run.font.size = DPt(10)
            run.font.name = 'Calibri'
            if bg_color:
                from docx.oxml.ns import qn
                shading = cell._element.get_or_add_tcPr()
                el = shading.makeelement(qn('w:shd'), {
                    qn('w:val'): 'clear', qn('w:color'): 'auto',
                    qn('w:fill'): bg_color})
                shading.append(el)
        return row

    def tabel_berjudul(headers, size=10):
        t = doc.add_table(rows=1, cols=len(headers))
        t.style = 'Table Grid'
        for i, h in enumerate(headers):
            t.rows[0].cells[i].text = h
            for p in t.rows[0].cells[i].paragraphs:
                p.runs[0].bold = True
                p.runs[0].font.size = DPt(size)
        return t

    # ── HALAMAN JUDUL ──
    for _ in range(4):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("PROPOSAL PENGEMBANGAN & IMPLEMENTASI")
    run.bold = True
    run.font.size = DPt(22)
    run.font.color.rgb = DRGBColor(0x0D, 0x94, 0x88)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"{NAMA_APLIKASI} — {NAMA_PANJANG.upper()}")
    run.bold = True
    run.font.size = DPt(20)
    run.font.color.rgb = DRGBColor(0x0F, 0x17, 0x2A)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Platform Siklus Penuh Pengelolaan Barang Milik Negara")
    run.font.size = DPt(15)
    run.font.color.rgb = DRGBColor(0x64, 0x74, 0x8B)
    for _ in range(2):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(
        "PWA offline-first · sinkron SIMAN V2 · persediaan standar SAKTI · "
        "tanda tangan elektronik · peta GIS\nSelaras SE-17/MK.1/2024, "
        "PP 27/2014 jo PP 28/2020, PMK 181/PMK.06/2016, PMK 234/PMK.05/2020")
    run.font.size = DPt(11)
    run.font.color.rgb = DRGBColor(0x64, 0x74, 0x8B)
    for _ in range(4):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"Versi Dokumen {VERSI} — Tahun {datetime.now().year}")
    run.bold = True
    run.font.size = DPt(13)
    doc.add_page_break()

    # ── DAFTAR ISI ──
    add_heading_styled("DAFTAR ISI", 1)
    toc = [
        ("BAB I", "PENDAHULUAN"), ("", "1.1 Latar Belakang"),
        ("", "1.2 Dasar Hukum"), ("", "1.3 Tujuan dan Sasaran"),
        ("", "1.4 Ruang Lingkup"),
        ("BAB II", "GAMBARAN UMUM SISTEM"), ("", "2.1 Deskripsi Sistem"),
        ("", "2.2 Arsitektur Teknologi"), ("", "2.3 Fitur dan Fungsionalitas"),
        ("BAB III", "METODOLOGI PENGEMBANGAN"),
        ("", "3.1 Metode Pengembangan"), ("", "3.2 Tahapan Implementasi"),
        ("", "3.3 Timeline"),
        ("BAB IV", "SPESIFIKASI TEKNIS"),
        ("", "4.1 Kebutuhan Perangkat Keras"),
        ("", "4.2 Kebutuhan Perangkat Lunak"),
        ("", "4.3 Kebutuhan Jaringan"),
        ("BAB V", "SKEMA LISENSI & RENCANA ANGGARAN BIAYA"),
        ("", "5.1 Skema Lisensi"), ("", "5.2 Referensi RAB Membangun Sendiri"),
        ("", "5.3 Biaya Operasional Tahunan"), ("", "5.4 Total Investasi"),
        ("BAB VI", "PENUTUP"),
    ]
    for bab, judul in toc:
        p = doc.add_paragraph()
        if bab:
            run = p.add_run(f"{bab}  ")
            run.bold = True
        run = p.add_run(judul)
        run.bold = bool(bab)
        p.paragraph_format.space_after = DPt(2)
    doc.add_page_break()

    # ── BAB I ──
    add_heading_styled("BAB I  PENDAHULUAN", 1)
    add_heading_styled("1.1 Latar Belakang", 2)
    add_para(
        "Barang Milik Negara (BMN) wajib dikelola akuntabel dan transparan "
        "sepanjang daur hidupnya — dari perencanaan kebutuhan sampai "
        "penghapusan. Praktik yang masih umum di satuan kerja: pencatatan "
        "tersebar di spreadsheet, foto dan dokumen terpisah dari data, "
        "rekonsiliasi dengan SIMAN/SAKTI dikerjakan manual, dan petugas "
        "lapangan terhambat karena aplikasi daring berhenti bekerja di "
        "lokasi tanpa sinyal.")
    add_para("Keterbatasan metode berjalan:", bold=True)
    for item in [
        "Rentan salah input, duplikasi, dan sulit diaudit (tanpa jejak per perubahan)",
        "Tidak real-time antar tim pencatat; rawan tumpang-tindih pencatatan",
        "Tidak berfungsi saat tidak ada koneksi — padahal inventarisasi fisik justru sering di titik tanpa sinyal",
        "Klasifikasi SE-17/MK.1/2024 (termasuk sub-klasifikasi aset tidak ditemukan) sulit dipenuhi manual",
        "Transaksi persediaan tidak mengikuti kode & FIFO per layer sebagaimana SAKTI",
        "Rekapitulasi, laporan resmi, dan berita acara memakan waktu berhari-hari",
    ]:
        p = doc.add_paragraph(item, style='List Bullet')
        p.paragraph_format.space_after = DPt(3)
    add_para(
        f"{NAMA_APLIKASI} ({NAMA_PANJANG}) dibangun menjawab kebutuhan itu: "
        "satu platform siklus penuh BMN — perencanaan, penganggaran, "
        "pengadaan, penggunaan, pemanfaatan, pengamanan & pemeliharaan, "
        "penilaian, pemindahtanganan, pemusnahan, penghapusan, penatausahaan "
        "(pembukuan–inventarisasi–pelaporan), dan pengawasan-pengendalian — "
        "yang tetap bekerja penuh saat sinyal hilang dan selaras dengan "
        "SIMAN V2 serta SAKTI.")

    add_heading_styled("1.2 Dasar Hukum", 2)
    for i, reg in enumerate([
        "Peraturan Pemerintah Nomor 27 Tahun 2014 jo PP Nomor 28 Tahun 2020 tentang Pengelolaan Barang Milik Negara/Daerah",
        "Peraturan Menteri Keuangan Nomor 181/PMK.06/2016 tentang Penatausahaan Barang Milik Negara",
        "Surat Edaran Menteri Keuangan Nomor SE-17/MK.1/2024 tentang Pelaksanaan Inventarisasi Barang Milik Negara",
        "Peraturan Menteri Keuangan Nomor 234/PMK.05/2020 tentang Kebijakan Akuntansi Pemerintah Pusat (persediaan metode FIFO)",
        "Peraturan Menteri Keuangan Nomor 207/PMK.06/2021 tentang Pengawasan dan Pengendalian BMN",
        "Peraturan ANRI Nomor 5 Tahun 2021 tentang Pedoman Umum Tata Naskah Dinas (persuratan)",
        "Peraturan Presiden Nomor 95 Tahun 2018 tentang Sistem Pemerintahan Berbasis Elektronik (SPBE)",
    ], 1):
        p = doc.add_paragraph(f"{i}. {reg}")
        p.paragraph_format.space_after = DPt(3)

    add_heading_styled("1.3 Tujuan dan Sasaran", 2)
    add_para("Tujuan:", bold=True)
    for obj in [
        "Menyediakan platform terintegrasi untuk seluruh siklus pengelolaan BMN dalam satu identitas aset",
        "Memenuhi klasifikasi SE-17/MK.1/2024 dan standar transaksi persediaan SAKTI (45 kode, FIFO per layer)",
        "Menjamin petugas lapangan tetap produktif tanpa koneksi (offline-first) dengan data aman tersinkron",
        "Mempercepat pelaporan: 15+ laporan resmi ber-kop siap tanda tangan, LBKP/LBP, dan CaLBMN pra-isi",
        "Meningkatkan akuntabilitas melalui jejak audit per field, jurnal Buku Barang, dan tanda tangan elektronik terverifikasi",
    ]:
        p = doc.add_paragraph(obj, style='List Bullet')
        p.paragraph_format.space_after = DPt(3)
    add_para("Sasaran:", bold=True)
    for target in [
        "100% data BMN dan persediaan tercatat digital, terstruktur, dan terisolasi per satuan kerja",
        "Pengurangan waktu inventarisasi dan penyusunan laporan hingga 60% dibanding metode manual",
        "Foto, dokumen sumber, BAST, dan nilai aset menyatu pada satu identitas aset",
        "Rekonsiliasi SIMAN V2 dan SAKTI dapat dilakukan langsung dari aplikasi",
    ]:
        p = doc.add_paragraph(target, style='List Bullet')
        p.paragraph_format.space_after = DPt(3)

    add_heading_styled("1.4 Ruang Lingkup", 2)
    add_para("Ruang lingkup sistem (kondisi terpasang saat ini):")
    for scope in [
        "Inventarisasi lapangan offline-first: mode lapangan satu layar, kamera ber-watermark GPS/jam, scan QR, antrean simpan persisten",
        "Manajemen aset 45+ field ber-registry (anti-drift) dengan OCC, audit trail, dan foto GridFS",
        "Siklus BMN 12 tahap: register & alur resmi per tahap + Wasdal (Lampiran PMK 207/2021)",
        "Penatausahaan: DBKP intra/ekstra, Buku Barang (jurnal SIMAK/SAKTI), KIB A–F, penyusutan per semester, riwayat nilai per aset",
        "Persediaan standar SAKTI: 45 kode transaksi, FIFO per layer, opname + BAOF, 7 laporan format SAKTI, daftar usang/rusak/tak dikuasai",
        "Pelaporan: 15+ laporan resmi PDF ber-kop, LBKP/LBP semesteran-tahunan, rekonsiliasi XLSX, CaLBMN pra-isi, batch unduh ZIP",
        "Sinkronisasi SIMAN V2 dua arah (impor/ekspor Master Aset, tandai selisih per aset, buat aset draft massal)",
        "Tanda tangan elektronik: kanvas/foto → PNG transparan, link per penanda tangan, QR + hash verifikasi, bubuhkan ke PDF unggahan",
        "Cetak stiker label BMN ber-QR: 3 ukuran × A4/A3, hierarki huruf terbaca, stiker contoh berukuran",
        "Spasial: peta aset Leaflet, impor/ekspor SHP/KML/KMZ/GeoJSON, denah dalam-gedung bergeoreferensi, geofence & pelacakan IoT ber-pagar privasi",
        "Master data: pegawai (impor massal + kartu UID e-KTP), pejabat, unit kerja Eselon I–V, satker & kop per satker, kodefikasi, akun BAS",
        "Persuratan: buku agenda & booking nomor naskah dinas lintas modul (PerANRI 5/2021)",
        "Keamanan: JWT + OTP email, RBAC 4 peran, isolasi multi-satker, rate-limit, backup otomatis terjadwal + restore + reset terlindungi",
    ]:
        p = doc.add_paragraph(scope, style='List Bullet')
        p.paragraph_format.space_after = DPt(3)
    doc.add_page_break()

    # ── BAB II ──
    add_heading_styled("BAB II  GAMBARAN UMUM SISTEM", 1)
    add_heading_styled("2.1 Deskripsi Sistem", 2)
    add_para(
        f"{NAMA_APLIKASI} adalah aplikasi web progresif (PWA) yang mencakup "
        "seluruh daur hidup BMN dalam satu identitas aset: transaksi setiap "
        "modul menulis jurnal ke Buku Barang, dokumen sumber menempel ke "
        "asetnya, dan laporan resmi dibangkitkan dari data yang sama. "
        "Pendekatan offline-first membuat petugas lapangan tetap bekerja "
        "penuh tanpa koneksi — snapshot kegiatan tersimpan di perangkat, "
        "antrean simpan tersinkron otomatis saat sinyal kembali.")
    add_para(
        "Kolaborasi multi-pengguna dijaga aman dengan optimistic concurrency "
        "control (version/If-Match), penguncian baris atomik, Idempotency-Key, "
        "dan WebSocket lintas worker; seluruh perubahan terekam pada jejak "
        "audit per field. Data antar satuan kerja terisolasi ketat — satu "
        "instans dapat melayani banyak satker dengan aman.")

    add_heading_styled("2.2 Arsitektur Teknologi", 2)
    tech_table = tabel_berjudul(["Lapisan", "Teknologi", "Deskripsi"])
    tech_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for row in [
        ("Presentasi (Frontend)", "React 19 + Tailwind CSS + shadcn/ui",
         "SPA/PWA offline-first: IndexedDB + Service Worker + antrean sinkron; Leaflet untuk peta & denah; tabel tervirtualisasi ribuan baris; mode gelap"),
        ("Logika Bisnis (Backend)", "Python FastAPI (async) + Motor",
         "45+ modul route; WebSocket + event bus lintas worker; OCC + Idempotency-Key; JWT + OTP email; isolasi multi-satker; ReportLab/WeasyPrint untuk PDF (di-offload ke thread)"),
        ("Data (Database)", "MongoDB 7 + GridFS",
         "Dokumen ber-UUID; GridFS untuk foto/dokumen/BAST; indeks kunci untuk daftar & filter; capped collection untuk event real-time"),
        ("Pendukung (Opsional)", "Redis + Meilisearch",
         "Cache bersama & rate-limit storage multi-worker (Redis); pencarian kilat toleran salah ketik (Meilisearch) — keduanya feature-flag, aplikasi tetap berjalan tanpanya"),
        ("Infrastruktur", "VPS: Nginx + Supervisor + GitHub Actions",
         "CI setiap PR (uji + lint + build), auto-deploy saat merge dengan gerbang health-check; SSL Let's Encrypt; backup otomatis terjadwal"),
    ]:
        add_table_row(tech_table, row)
    doc.add_paragraph()
    add_para("Keunggulan arsitektur:", bold=True)
    for adv in [
        "Offline-first sungguhan: snapshot + antrean persisten, bukan sekadar cache halaman",
        "Aman multi-pengguna: OCC, locking, idempotensi — tanpa lost-update dan duplikasi",
        "Registry terpusat (field aset, kode transaksi SAKTI) + test anti-drift menjaga konsistensi lintas modul",
        "Semua teks antarmuka, laporan, dan dokumen berbahasa Indonesia",
    ]:
        p = doc.add_paragraph(adv, style='List Bullet')
        p.paragraph_format.space_after = DPt(3)

    add_heading_styled("2.3 Fitur dan Fungsionalitas", 2)
    for judul, butir in [
        ("A. Inventarisasi Lapangan (Offline-First)", [
            "Mode lapangan satu layar: status & kondisi sekali ketuk, salin lokasi/pengguna aset sebelumnya, GPS cache instan",
            "Kamera penuh ala Timemark: watermark jam & koordinat, flash, gestur kecerahan, alur beruntun Simpan & Baru / Simpan & Scan QR",
            "Snapshot kegiatan di IndexedDB (delta sync) + antrean simpan persisten — bekerja penuh tanpa koneksi",
            "Klasifikasi SE-17/MK.1/2024 lengkap: 5 status, sub-klasifikasi Tidak Ditemukan, progres % real-time"]),
        ("B. Manajemen Aset & Media", [
            "45+ field per aset dari satu registry (anti-drift): identitas, perolehan, organisasi, kondisi, GPS, stiker, dokumen, garansi",
            "Multi-foto GridFS berkompresi berlapis + streaming cacheable (ETag/304); checklist dokumen ber-bukti",
            "Barang Serupa (grup per kode) + Ubah Massal ber-jurnal; impor/ekspor 46 kolom CSV/XLSX dua arah dengan template dropdown"]),
        ("C. Siklus BMN 12 Tahap", [
            "Perencanaan (RKBMN + sanding SBSK), Penganggaran, Pengadaan (BAST/LPB + pencatatan otomatis), Penggunaan (PSP, BAST pengguna)",
            "Pemanfaatan, Pengamanan & Pemeliharaan (BA Perbaikan + dampak masa manfaat), Penilaian (revaluasi ber-jurnal)",
            "Pemindahtanganan, Pemusnahan, Penghapusan (SK + register), Wasdal (Lampiran PMK 207/2021 + investigasi)"]),
        ("D. Penatausahaan & Nilai", [
            "DBKP per golongan intra/ekstra + Buku Barang: jurnal mutasi ber-kode SIMAK/SAKTI append-only dari SEMUA modul",
            "KIB A–F per unit (PMK 181) ber-foto & riwayat; penyusutan per semester selaras SIMAN",
            "Riwayat Nilai per aset: perolehan → kapitalisasi/revaluasi/koreksi → nilai buku; edit harga manual/massal pun berjurnal"]),
        ("E. Persediaan Standar SAKTI", [
            "45 kode transaksi resmi (M01–M99, K01–K99, P01, H01–H03) dengan arah eksplisit + Daftar Transaksi lintas barang ber-filter",
            "FIFO per layer (PMK 234/2020) + koreksi nilai proporsional; opname semesteran + kertas kerja + BAOF 3 penandatangan",
            "Daftar Usang/Rusak/Tak Dikuasai (bahan CaLK) + hapus definitif ber-SK; 7 laporan PDF persis format SAKTI"]),
        ("F. Pelaporan Resmi", [
            "15+ laporan PDF ber-kop: DBHI 8 tipe, RHI, BAHI, BA, SPTJM, Surat Koreksi, Daftar Pemegang, Eksekutif, LHI, LBP lengkap",
            "LBKP semesteran/tahunan + rekonsiliasi XLSX + CaLBMN pra-isi; batch unduh ZIP seluruh laporan per kegiatan",
            "Pengesahan berkekuatan dokumen: unggah PDF bertanda tangan → kegiatan terkunci permanen + kartu riwayat per aset"]),
        ("G. Tanda Tangan Elektronik", [
            "Spesimen dari kanvas goresan mulus atau foto kertas → PNG transparan",
            "Permintaan TTD via link per penanda tangan (token sekali-pakai, urutan, kedaluwarsa 14 hari) + notifikasi email",
            "Bubuhkan TTD + QR verifikasi ke PDF laporan/BAST/dokumen unggahan — atur letak & ukuran di pratinjau; verifikasi publik ber-hash (NIP di-masking)"]),
        ("H. Stiker Label & Kartu", [
            "Stiker QR 3 ukuran × A4/A3 penuh-halaman: hierarki huruf terbaca di semua ukuran, nama panjang lanjut baris, stiker contoh berukuran",
            "Kartu BMN format KTP (satuan & massal) + Kartu Pegawai UID e-KTP/NFC untuk isi identitas cepat",
            "Scan QR kamera: buka aset, catat observasi lokasi, rekonsiliasi opname ruangan"]),
        ("I. Spasial & Pelacakan", [
            "Peta aset interaktif mengikuti filter aktif (offline-ready) + ekspor KML/KMZ/SHP ber-atribut",
            "Hierarki spasial + denah dalam-gedung bergeoreferensi (impor SHP/KML/KMZ/GeoJSON, gambar poligon, belah wilayah)",
            "Geofence histeresis, pelacakan aset bergerak & IoT ber-pagar privasi (UU PDP), SBSK dari luas poligon nyata"]),
        ("J. SIMAN V2 & Integrasi", [
            "Impor/ekspor 'Master Aset' SIMAN V2 — deteksi header semua sheet, unggah andal (retry), validasi satker 6↔20 digit",
            "Tandai selisih per aset (≠ SIMAN) + sinkron 1-klik; baris belum tercatat → CSV atau buat aset draft massal",
            "Persuratan terpadu: buku agenda + booking nomor lintas modul (PerANRI 5/2021)"]),
        ("K. Master Data & SDM", [
            "Master Pegawai (impor Excel massal, foto, validasi rekening/WNI-WNA, tanggalan roda Indonesia), Pejabat, Unit Kerja Eselon I–V",
            "Master Satker + kop per satker; Kodefikasi barang; Referensi Akun BAS per makna digit (KEP-211/PB/2018)",
            "Keterkaitan aset ↔ pegawai: panel Perlu Serah Terima BMN + peringatan status kepegawaian"]),
        ("L. Keamanan, Audit & Data", [
            "JWT + OTP email, lockout login, revokasi sesi; RBAC 4 peran (super admin, admin, operator, viewer read-only ditegakkan server)",
            "Isolasi multi-satker menyeluruh (baca/tulis/hapus/ekspor); jejak audit per field + log sistem",
            "Backup otomatis terjadwal harian + retensi + restore dari arsip; reset melindungi master referensi; anti path-traversal"]),
    ]:
        add_para(judul, bold=True, space_after=3)
        for item in butir:
            p = doc.add_paragraph(item, style='List Bullet')
            p.paragraph_format.space_after = DPt(2)
        doc.add_paragraph()
    doc.add_page_break()

    # ── BAB III ──
    add_heading_styled("BAB III  METODOLOGI PENGEMBANGAN", 1)
    add_heading_styled("3.1 Metode Pengembangan", 2)
    add_para(
        "Pengembangan memakai metodologi Agile inkremental: satu fitur satu "
        "pull request, setiap perubahan melewati gerbang otomatis (uji unit "
        "backend, lint, build produksi) sebelum digabung, dan penggabungan "
        "memicu deploy otomatis dengan pemeriksaan kesehatan pasca-restart. "
        "Pendekatan ini memungkinkan rilis kecil yang sering, mudah "
        "di-rollback, dan selalu terdokumentasi pada CHANGELOG.")
    add_para("Praktik mutu yang dijalankan:", bold=True)
    for b in [
        "1.400+ uji unit backend bebas-infrastruktur dijalankan pada setiap perubahan",
        "Registry terpusat + test anti-drift (field aset, kode transaksi SAKTI, template impor/ekspor)",
        "Tinjauan adversarial berkala: temuan keamanan/keandalan ditutup bergelombang dan terdokumentasi",
        "Smoke-test laporan PDF dengan merender dokumen sungguhan tanpa database",
    ]:
        p = doc.add_paragraph(b, style='List Bullet')

    add_heading_styled("3.2 Tahapan Implementasi", 2)
    for judul, butir in [
        ("Fase 1: Analisis & Penyesuaian (Bulan 1-2)", [
            "Pemetaan struktur satker, kodefikasi, dan kop laporan instansi",
            "Penyesuaian master data (pegawai, pejabat, unit kerja, akun BAS)",
            "Penetapan kebijakan: ambang kapitalisasi, jadwal opname, peran pengguna"]),
        ("Fase 2: Instalasi & Migrasi (Bulan 3-5)", [
            "Provisioning VPS/on-premise + domain + SSL",
            "Migrasi data awal: impor SIMAN V2 / CSV / XLSX + validasi",
            "Konfigurasi backup otomatis terjadwal + uji pulihkan"]),
        ("Fase 3: Uji Terima & Pelatihan (Bulan 6-7)", [
            "UAT per modul bersama perwakilan pengguna",
            "Pelatihan admin & operator (termasuk mode offline lapangan)",
            "Gladi inventarisasi ujung-ke-ujung sampai pengesahan"]),
        ("Fase 4: Go-Live (Bulan 8)", [
            "Soft launch + pendampingan intensif",
            "Serah terima dokumentasi & kredensial",
            "Penetapan kanal dukungan"]),
        ("Fase 5: Pemeliharaan (Bulan 9-12)", [
            "Pemantauan, perbaikan, dan pembaruan berkala",
            "Evaluasi pasca-kegiatan inventarisasi pertama",
            "Perencanaan pengembangan lanjutan"]),
    ]:
        add_para(judul, bold=True, space_after=3)
        for item in butir:
            p = doc.add_paragraph(item, style='List Bullet')
            p.paragraph_format.space_after = DPt(2)
        doc.add_paragraph()

    add_heading_styled("3.3 Timeline", 2)
    tl = tabel_berjudul(["Fase", "Kegiatan", "Bln 1-2", "Bln 3-5", "Bln 6-7",
                         "Bln 8", "Bln 9-12"], size=9)
    for row in [
        ("1", "Analisis & Penyesuaian", "████", "", "", "", ""),
        ("2", "Instalasi & Migrasi", "", "████", "", "", ""),
        ("3", "Uji Terima & Pelatihan", "", "", "████", "", ""),
        ("4", "Go-Live", "", "", "", "████", ""),
        ("5", "Pemeliharaan", "", "", "", "", "████"),
    ]:
        add_table_row(tl, row)
    doc.add_page_break()

    # ── BAB IV ──
    add_heading_styled("BAB IV  SPESIFIKASI TEKNIS", 1)
    add_heading_styled("4.1 Kebutuhan Perangkat Keras (Server)", 2)
    hw = tabel_berjudul(["Komponen", "Minimum", "Rekomendasi", "Keterangan"])
    for row in [
        ("Processor", "4 vCPU", "8 vCPU", "Multi-worker + render PDF"),
        ("Memory (RAM)", "8 GB", "16 GB", "PDF berat + cache"),
        ("Storage (SSD)", "100 GB", "500 GB", "Foto GridFS & arsip backup"),
        ("Bandwidth", "100 Mbps", "1 Gbps", "Unggah foto & sinkron lapangan"),
        ("OS", "Ubuntu 22.04 LTS", "Ubuntu 24.04 LTS", "Server Linux"),
    ]:
        add_table_row(hw, row)

    add_heading_styled("4.2 Kebutuhan Perangkat Lunak", 2)
    sw = tabel_berjudul(["Software", "Versi", "Fungsi", "Lisensi"])
    for row in [
        ("Python", "3.11+", "Runtime backend (FastAPI)", "Open Source"),
        ("Node.js", "20 LTS", "Build frontend", "Open Source"),
        ("MongoDB", "7.0+", "Database + GridFS", "SSPL (gratis)"),
        ("Nginx", "1.24+", "Reverse proxy & SSL", "Open Source"),
        ("Supervisor", "4+", "Manajemen proses", "Open Source"),
        ("React", "19+", "Framework frontend (PWA)", "MIT"),
        ("FastAPI", "0.110+", "Framework backend", "MIT"),
        ("Redis (opsional)", "7+", "Cache & rate-limit bersama", "Open Source"),
        ("Meilisearch (opsional)", "1.x", "Pencarian kilat", "MIT"),
    ]:
        add_table_row(sw, row)

    add_heading_styled("4.3 Kebutuhan Jaringan", 2)
    for req in [
        "Koneksi internet stabil untuk server (klien lapangan boleh offline — sinkron saat sinyal kembali)",
        "Alamat IP publik statis + domain dengan sertifikat SSL (Let's Encrypt atau komersial)",
        "Firewall: port 80 (redirect), 443 (HTTPS); MongoDB/Redis/Meilisearch hanya internal",
        "Dukungan WebSocket pada reverse proxy",
        "Layanan email SMTP/API (OTP & notifikasi tanda tangan elektronik)",
    ]:
        p = doc.add_paragraph(req, style='List Bullet')
    doc.add_page_break()

    # ── BAB V ──
    add_heading_styled("BAB V  SKEMA LISENSI & RENCANA ANGGARAN BIAYA", 1)
    add_heading_styled("5.1 Skema Lisensi", 2)
    add_para(
        f"{NAMA_APLIKASI} tersedia sebagai produk berlisensi — instansi "
        "tidak menanggung risiko membangun dari nol:")
    lic = tabel_berjudul(["Paket", "Harga", "Cakupan"])
    for nama, harga, tagline, fitur_l in TIER_LISENSI:
        add_table_row(lic, (nama, harga, tagline + ". " + "; ".join(fitur_l)))
    doc.add_paragraph()
    add_para(
        "Harga di atas adalah acuan penawaran; nilai final mengikuti "
        "negosiasi, lingkup kustomisasi, dan ketentuan pengadaan yang "
        "berlaku. Tanpa biaya per-user — seluruh tim dapat bekerja "
        "bersamaan.", italic=True)

    add_heading_styled("5.2 Referensi RAB — Membangun & Mengoperasikan Sendiri", 2)
    add_para(
        "Sebagai pembanding kewajaran harga, berikut estimasi biaya bila "
        "instansi membangun sistem setara dari nol:")
    rab = tabel_berjudul(["No", "Komponen Biaya", "Volume", "Satuan",
                          "Harga Satuan (Rp)", "Jumlah (Rp)"], size=9)
    rab.alignment = WD_TABLE_ALIGNMENT.CENTER
    for row in RAB_BARIS:
        is_header = row[0] in ("A", "B", "C") or row[0] == ""
        add_table_row(rab, row, bold=is_header)
    add_table_row(rab, ["", "TOTAL TAHUN PERTAMA", "", "", "", RAB_TOTAL_TH1],
                  bold=True, bg_color="0D9488")
    doc.add_paragraph()

    add_heading_styled("5.3 Biaya Operasional Tahunan (Tahun ke-2 dst)", 2)
    annual = tabel_berjudul(["No", "Komponen", "Biaya/Bulan (Rp)",
                             "Biaya/Tahun (Rp)"])
    for row in [
        ("1", "Cloud Server (VPS)", "1.500.000", "18.000.000"),
        ("2", "MongoDB Atlas", "800.000", "9.600.000"),
        ("3", "Domain & SSL", "-", "500.000"),
        ("4", "Backup Storage", "200.000", "2.400.000"),
        ("5", "Maintenance & Support", "3.000.000", "36.000.000"),
    ]:
        add_table_row(annual, row)
    add_table_row(annual, ["", "TOTAL PER TAHUN", "", RAB_TAHUNAN], bold=True,
                  bg_color="0D9488")

    add_heading_styled("5.4 Total Investasi", 2)
    inv = tabel_berjudul(["Periode", "Komponen", "Biaya (Rp)"])
    for row in [
        ("Tahun 1", "Pengembangan + Infrastruktur + Operasional", RAB_TOTAL_TH1),
        ("Tahun 2", "Infrastruktur + Operasional", RAB_TAHUNAN),
        ("Tahun 3", "Infrastruktur + Operasional", RAB_TAHUNAN),
    ]:
        add_table_row(inv, row)
    add_table_row(inv, ["", "TOTAL INVESTASI 3 TAHUN (membangun sendiri)",
                        RAB_TOTAL_3TH], bold=True, bg_color="0D9488")
    doc.add_page_break()

    # ── BAB VI ──
    add_heading_styled("BAB VI  PENUTUP", 1)
    add_heading_styled("6.1 Kesimpulan", 2)
    add_para(
        f"{NAMA_APLIKASI} ({NAMA_PANJANG}) {VERSI} adalah platform siklus "
        "penuh pengelolaan BMN yang sudah berjalan dan teruji — bukan "
        "rencana pengembangan. Dengan offline-first untuk lapangan, "
        "keselarasan SIMAN V2 & SAKTI, tanda tangan elektronik terverifikasi, "
        "dan 15+ laporan resmi siap sah, sistem ini:")
    for item in [
        "Memangkas waktu inventarisasi dan pelaporan hingga 60% dibanding metode manual",
        "Menjamin akurasi melalui registry terpusat, OCC multi-user, jurnal Buku Barang, dan jejak audit per field",
        "Memenuhi SE-17/MK.1/2024, PMK 181/2016, PMK 234/2020, dan PMK 207/2021 dalam satu aplikasi",
        "Melayani banyak satuan kerja dalam satu instans dengan isolasi data yang ketat",
        "Menjaga keberlangsungan data lewat backup otomatis terjadwal dan pemulihan dari arsip",
    ]:
        p = doc.add_paragraph(item, style='List Bullet')
    add_heading_styled("6.2 Rekomendasi", 2)
    for item in [
        "Memilih skema lisensi sesuai cakupan satker (tunggal / multi-satker / perpetual + source code)",
        "Menetapkan tim pengelola: administrator sistem + penanggung jawab BMN per satker",
        "Menjadwalkan migrasi data (SIMAN V2/CSV) dan pelatihan sebelum kegiatan inventarisasi berikutnya",
        "Memanfaatkan backup otomatis + cadangan infrastruktur (mongodump terjadwal) sejak hari pertama",
        "Evaluasi berkala pasca-kegiatan untuk penyesuaian kebijakan dan fitur",
    ]:
        p = doc.add_paragraph(item, style='List Bullet')
    doc.add_paragraph()
    add_para(
        "Demikian proposal ini disampaikan sebagai bahan pertimbangan. Atas "
        "perhatian dan persetujuannya, kami ucapkan terima kasih.")
    for _ in range(3):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    p.add_run(f"Nusantara, {datetime.now().strftime('%d %B %Y')}")
    for _ in range(3):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = p.add_run(f"Tim Pengembang\n{NAMA_APLIKASI} — {NAMA_PANJANG}")
    run.bold = True

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument."
                   "wordprocessingml.document",
        headers={"Content-Disposition":
                 "attachment; filename=Proposal_AMAN.docx"})
